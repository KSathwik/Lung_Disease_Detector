"""
Generates the LungAI project presentation (PPTX).
Run: python docs/generate_ppt.py
Output: docs/LungAI_Presentation.pptx
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── Theme ────────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x0B, 0x1F, 0x3A)
BLUE   = RGBColor(0x18, 0x5F, 0xA5)
TEAL   = RGBColor(0x0F, 0x6E, 0x56)
LIGHT  = RGBColor(0xF2, 0xF6, 0xFB)
GREY   = RGBColor(0x55, 0x60, 0x6E)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x2E, 0xA3, 0xD9)

SW, SH = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def rect(s, x, y, w, h, color, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=6, line_spacing=1.05):
    """runs: list of paragraphs; each paragraph is list of (txt, size, color, bold)."""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        for (txt, size, color, bold) in para:
            r = p.add_run()
            r.text = txt
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.name = "Calibri"
    return tb


def header(s, title, kicker=None):
    rect(s, 0, 0, SW, Inches(1.15), NAVY)
    rect(s, 0, Inches(1.15), SW, Inches(0.06), ACCENT)
    text(s, Inches(0.6), Inches(0.18), Inches(12), Inches(0.8),
         [[(title, 30, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
    if kicker:
        text(s, Inches(0.62), Inches(0.74), Inches(12), Inches(0.35),
             [[(kicker, 12, ACCENT, True)]])


def bullets(s, x, y, w, h, items, size=15, color=NAVY, gap=8):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.line_spacing = 1.05
        if isinstance(it, tuple):
            label, desc = it
            r = p.add_run(); r.text = "•  " + label
            r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = color; r.font.name = "Calibri"
            r2 = p.add_run(); r2.text = "  —  " + desc
            r2.font.size = Pt(size); r2.font.color.rgb = GREY; r2.font.name = "Calibri"
        else:
            r = p.add_run(); r.text = "•  " + it
            r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = "Calibri"
    return tb


def card(s, x, y, w, h, title, body_items, accent=BLUE):
    rect(s, x, y, w, h, LIGHT)
    rect(s, x, y, w, Inches(0.5), accent)
    text(s, x + Inches(0.15), y + Inches(0.04), w - Inches(0.3), Inches(0.42),
         [[(title, 15, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
    bullets(s, x + Inches(0.2), y + Inches(0.62), w - Inches(0.4), h - Inches(0.7),
            body_items, size=12.5, gap=5)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(4.55), SW, Inches(0.07), ACCENT)
text(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(1.2),
     [[("LungAI", 60, WHITE, True)]])
text(s, Inches(0.95), Inches(3.15), Inches(11.5), Inches(0.9),
     [[("AI-Powered Lung Disease Detection from Chest X-rays", 26, ACCENT, True)]])
text(s, Inches(0.95), Inches(4.75), Inches(11.5), Inches(1.0),
     [[("Deep Learning system comparing a custom CNN and ResNet50 transfer learning", 16, RGBColor(0xC9,0xD6,0xE6), False)],
      [("FastAPI  •  React  •  TensorFlow / Keras", 15, RGBColor(0x8F,0xA6,0xC4), True)]],
     space_after=8)
text(s, Inches(0.95), Inches(6.6), Inches(11.5), Inches(0.5),
     [[("Academic Project  |  v1.0", 13, RGBColor(0x6E,0x86,0xA6), True)]])

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Problem & Motivation
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
header(s, "Problem & Motivation", "WHY THIS MATTERS")
bullets(s, Inches(0.6), Inches(1.5), Inches(6.0), Inches(5.5),
        [("Burden", "Lung diseases (pneumonia, TB, COVID-19) are leading causes of death worldwide."),
         ("Bottleneck", "Chest X-ray interpretation needs trained radiologists, who are scarce in many regions."),
         ("Delay", "Slow diagnosis delays treatment and worsens outcomes."),
         ("Variability", "Manual reading is subjective and error-prone under heavy workloads."),
         ("Opportunity", "Deep learning can flag abnormalities quickly and consistently as a decision-support aid.")],
        size=15, gap=12)
card(s, Inches(6.95), Inches(1.55), Inches(5.75), Inches(4.4),
     "Project Goal",
     ["Detect 7 lung conditions from a single chest X-ray",
      "Train & compare two algorithms (CNN vs ResNet50)",
      "Automatically select the best-performing model",
      "Serve predictions through a clean web application",
      "Store patients, scans, predictions & reports"],
     accent=TEAL)
text(s, Inches(0.6), Inches(6.75), Inches(12), Inches(0.5),
     [[("Disclaimer: An academic decision-support prototype — it assists, not replaces, qualified physicians.", 12, GREY, True)]])

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — System Architecture
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
header(s, "System Architecture", "END-TO-END PIPELINE")
card(s, Inches(0.5), Inches(1.55), Inches(3.9), Inches(4.6), "Frontend  (React)",
     ["Analyze: upload X-ray & view result",
      "Metrics: CNN vs ResNet charts",
      "History: past predictions",
      "Patients: record management",
      "Axios API service layer"], accent=ACCENT)
card(s, Inches(4.7), Inches(1.55), Inches(3.9), Inches(4.6), "Backend  (FastAPI)",
     ["REST API under /api/v1",
      "predict / predictions endpoints",
      "patients & reports CRUD",
      "ML inference engine",
      "Health check + logging"], accent=BLUE)
card(s, Inches(8.9), Inches(1.55), Inches(3.9), Inches(4.6), "Data & ML",
     ["TensorFlow / Keras models",
      "Preprocessing pipeline",
      "SQLAlchemy ORM database",
      "Trained .h5 model files",
      "Training curves & matrices"], accent=TEAL)
text(s, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.7),
     [[("Flow:  ", 14, NAVY, True),
       ("User uploads X-ray  →  React  →  FastAPI  →  Preprocess  →  Model inference  →  Prediction stored  →  Result + report",
        14, GREY, False)]])

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Data Preprocessing Pipeline
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
header(s, "Data Preprocessing Pipeline", "FROM RAW X-RAYS TO MODEL-READY TENSORS")
steps = [("1. Scan", "Index all images & labels"),
         ("2. Clean", "Drop duplicates, corrupt & tiny images"),
         ("3. Enhance", "CLAHE contrast (critical for X-rays)"),
         ("4. Denoise", "Gaussian blur"),
         ("5. Split", "70% train / 15% val / 15% test (stratified)"),
         ("6. Augment", "Flip, rotate ±15°, brightness jitter (train only)"),
         ("7. Normalize", "ImageNet mean / std normalization")]
y = Inches(1.55)
for i, (t, d) in enumerate(steps):
    col = i % 2
    row = i // 2
    x = Inches(0.6) + col * Inches(6.3)
    yy = y + row * Inches(0.92)
    rect(s, x, yy, Inches(6.0), Inches(0.78), LIGHT)
    rect(s, x, yy, Inches(0.12), Inches(0.78), ACCENT if col == 0 else TEAL)
    text(s, x + Inches(0.28), yy + Inches(0.04), Inches(5.6), Inches(0.7),
         [[(t + "   ", 15, NAVY, True), (d, 13, GREY, False)]], anchor=MSO_ANCHOR.MIDDLE)
text(s, Inches(0.6), Inches(6.75), Inches(12), Inches(0.5),
     [[("Output: 224×224×3 tensors  •  7 balanced classes  •  augmentation applied only to the training split", 13, GREY, True)]])

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Diseases Detected
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
header(s, "Conditions Detected", "7-CLASS CLASSIFICATION")
classes = ["Normal", "Pneumonia", "Tuberculosis", "COVID-19",
           "Lung Cancer", "COPD", "Pleural Effusion"]
colors = [TEAL, BLUE, RGBColor(0x8A,0x4F,0x9E), RGBColor(0xB0,0x3A,0x3A),
          RGBColor(0x6E,0x2C,0x2C), RGBColor(0x2C,0x5F,0x6E), RGBColor(0x4F,0x6E,0x2C)]
for i, (c, col) in enumerate(zip(classes, colors)):
    cx = i % 4
    cy = i // 4
    x = Inches(0.6) + cx * Inches(3.15)
    yy = Inches(1.7) + cy * Inches(1.7)
    rect(s, x, yy, Inches(2.9), Inches(1.45), col)
    text(s, x + Inches(0.15), yy + Inches(0.1), Inches(2.6), Inches(1.25),
         [[(c, 18, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
text(s, Inches(0.6), Inches(6.6), Inches(12), Inches(0.6),
     [[("Softmax output over 7 classes  •  input 224×224×3  •  per-class confidence returned with each prediction", 13, GREY, True)]])

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Algorithm 1: Custom CNN
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
header(s, "Algorithm 1 — Custom CNN", "BUILT FROM SCRATCH")
bullets(s, Inches(0.6), Inches(1.55), Inches(6.0), Inches(5),
        [("Type", "Custom convolutional network, trained end-to-end"),
         ("Depth", "4 convolutional blocks (32 → 64 → 128 → 256 filters)"),
         ("Each block", "Conv2D + BatchNorm + ReLU + MaxPool + Dropout"),
         ("Head", "GlobalAvgPool → Dense(512) → Dropout(0.5) → Softmax"),
         ("Optimizer", "Adam, learning rate 1e-4"),
         ("Loss", "Sparse categorical cross-entropy"),
         ("Best for", "Smaller, focused datasets")],
        size=14, gap=11)
card(s, Inches(6.95), Inches(1.55), Inches(5.75), Inches(4.6), "Layer Flow",
     ["Input 224×224×3",
      "Conv Block 1  →  32 filters",
      "Conv Block 2  →  64 filters",
      "Conv Block 3  →  128 filters",
      "Conv Block 4  →  256 filters + GAP",
      "Dense 512 + Dropout 0.5",
      "Dense 7 + Softmax"], accent=BLUE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Algorithm 2: ResNet50
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
header(s, "Algorithm 2 — ResNet50 Transfer Learning", "PRETRAINED ON IMAGENET")
bullets(s, Inches(0.6), Inches(1.55), Inches(6.0), Inches(5),
        [("Base", "ResNet50, 50-layer residual network"),
         ("Weights", "Pretrained on ImageNet"),
         ("Head", "GAP → Dense(1024) → Dense(512) → Softmax"),
         ("Phase 1", "Freeze base, train classifier head (10 epochs, LR 1e-3)"),
         ("Phase 2", "Unfreeze top 30 layers, fine-tune (LR 1e-5)"),
         ("Best for", "Larger, diverse datasets")],
        size=14, gap=12)
card(s, Inches(6.95), Inches(1.55), Inches(5.75), Inches(4.6), "Two-Phase Strategy",
     ["Load ResNet50 (ImageNet, no top)",
      "Phase 1: base frozen → learn head",
      "Phase 2: unfreeze top 30 layers",
      "Fine-tune with very low LR",
      "Callbacks: EarlyStopping,",
      "ReduceLROnPlateau, Checkpoint"], accent=TEAL)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Model Comparison & Selection
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
header(s, "Model Comparison & Selection", "AUTOMATIC BEST-MODEL PICK")
# table-ish comparison
rows = [("", "Custom CNN", "ResNet50"),
        ("Approach", "From scratch", "Transfer learning"),
        ("Depth", "4 conv blocks", "50 residual layers"),
        ("Pretrained", "No", "ImageNet"),
        ("Training", "End-to-end", "2-phase fine-tune"),
        ("Strength", "Small datasets", "Large datasets")]
tx, ty = Inches(0.6), Inches(1.55)
cw = [Inches(3.0), Inches(4.4), Inches(4.4)]
for r, row in enumerate(rows):
    x = tx
    for c, val in enumerate(row):
        head_row = (r == 0)
        fill = NAVY if head_row else (LIGHT if r % 2 else WHITE)
        rect(s, x, ty + r*Inches(0.6), cw[c], Inches(0.6), fill,
             line=RGBColor(0xD0,0xD8,0xE2))
        col = WHITE if head_row else NAVY
        bold = head_row or c == 0
        text(s, x + Inches(0.12), ty + r*Inches(0.6), cw[c]-Inches(0.2), Inches(0.6),
             [[(val, 13.5, col, bold)]], anchor=MSO_ANCHOR.MIDDLE)
        x += cw[c]
card(s, Inches(0.6), Inches(5.5), Inches(12.1), Inches(1.55), "Selection Formula",
     ["Composite score  =  0.4 × F1  +  0.3 × Accuracy  +  0.2 × AUC-ROC  +  0.1 × Recall   →   higher score wins (ResNet50 typically on larger data)"],
     accent=ACCENT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Evaluation Metrics
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
header(s, "Evaluation & Metrics", "HOW WE MEASURE QUALITY")
metrics = [("Accuracy", "Overall correct predictions"),
           ("Precision", "Correctness of positive calls"),
           ("Recall", "Coverage of actual positives"),
           ("F1-Score", "Balance of precision & recall"),
           ("AUC-ROC", "Class separability (macro OvR)"),
           ("Confusion Matrix", "Per-class error breakdown")]
for i, (t, d) in enumerate(metrics):
    cx = i % 3
    cy = i // 3
    x = Inches(0.6) + cx * Inches(4.15)
    yy = Inches(1.7) + cy * Inches(1.85)
    rect(s, x, yy, Inches(3.85), Inches(1.55), LIGHT)
    rect(s, x, yy, Inches(3.85), Inches(0.12), BLUE if cy == 0 else TEAL)
    text(s, x + Inches(0.2), yy + Inches(0.22), Inches(3.5), Inches(1.2),
         [[(t, 17, NAVY, True)], [(d, 12.5, GREY, False)]], space_after=6)
text(s, Inches(0.6), Inches(6.7), Inches(12), Inches(0.5),
     [[("Generated artifacts: training_results.json, training_curves.png, confusion_matrix_CNN.png, confusion_matrix_ResNet.png", 12, GREY, True)]])

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Tech Stack
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
header(s, "Technology Stack", "TOOLS & FRAMEWORKS")
card(s, Inches(0.5), Inches(1.55), Inches(3.9), Inches(4.6), "Frontend",
     ["React (SPA)", "React Router", "Axios", "Chart components", "Custom CSS"], accent=ACCENT)
card(s, Inches(4.7), Inches(1.55), Inches(3.9), Inches(4.6), "Backend",
     ["FastAPI", "Uvicorn", "SQLAlchemy (async ORM)", "Pydantic", "Python 3.10+"], accent=BLUE)
card(s, Inches(8.9), Inches(1.55), Inches(3.9), Inches(4.6), "ML & Data",
     ["TensorFlow / Keras", "ResNet50 (ImageNet)", "scikit-learn", "OpenCV (CLAHE)", "NumPy, Matplotlib, Seaborn"], accent=TEAL)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — API & Database
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
header(s, "API & Database", "INTERFACES AND STORAGE")
bullets(s, Inches(0.6), Inches(1.55), Inches(6.0), Inches(5),
        ["GET  /api/v1/health  —  health check",
         "POST /api/v1/predict  —  upload image, get prediction",
         "GET  /api/v1/predictions  —  list all predictions",
         "GET  /api/v1/predictions/{id}  —  single prediction",
         "GET  /api/v1/model-metrics  —  CNN vs ResNet metrics",
         "POST /api/v1/patients  —  register patient",
         "POST /api/v1/reports/generate/{id}  —  report"],
        size=13.5, gap=10)
card(s, Inches(6.95), Inches(1.55), Inches(5.75), Inches(4.6), "Database Schema",
     ["patients — demographics & history",
      "lung_scans — uploaded image metadata",
      "predictions — CNN + ResNet results",
      "model_metrics — per-version metrics",
      "reports — generated reports"], accent=BLUE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Workflow / Demo
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
header(s, "Application Workflow", "TYPICAL USER JOURNEY")
flow = ["Register / select a patient",
        "Upload a chest X-ray image",
        "Backend preprocesses the image (resize, CLAHE, normalize)",
        "Selected model runs inference → 7-class probabilities",
        "Top prediction + confidence shown in the UI",
        "Result saved to history; report can be generated"]
for i, step in enumerate(flow):
    yy = Inches(1.6) + i * Inches(0.82)
    rect(s, Inches(0.6), yy, Inches(0.7), Inches(0.62), BLUE if i % 2 == 0 else TEAL)
    text(s, Inches(0.6), yy, Inches(0.7), Inches(0.62),
         [[(str(i+1), 20, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(1.5), yy, Inches(11), Inches(0.62),
         [[(step, 16, NAVY, False)]], anchor=MSO_ANCHOR.MIDDLE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Challenges & Future Work
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
header(s, "Challenges & Future Work", "LIMITATIONS AND NEXT STEPS")
card(s, Inches(0.6), Inches(1.55), Inches(5.9), Inches(4.7), "Challenges",
     ["Data quality & class imbalance",
      "Limited labeled medical images",
      "Risk of overfitting on small data",
      "Compute cost of training ResNet50",
      "Generalizing across X-ray machines"], accent=RGBColor(0xB0,0x3A,0x3A))
card(s, Inches(6.85), Inches(1.55), Inches(5.9), Inches(4.7), "Future Work",
     ["DICOM support via pydicom",
      "Grad-CAM explainability heatmaps",
      "GPU deployment (NVIDIA T4+)",
      "Larger datasets (5,000+ / class)",
      "Clinical validation & PDF reports"], accent=TEAL)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Conclusion
# ══════════════════════════════════════════════════════════════════════════════
s = slide()
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(2.35), SW, Inches(0.07), ACCENT)
text(s, Inches(0.9), Inches(0.9), Inches(11.5), Inches(1.3),
     [[("Conclusion", 44, WHITE, True)]])
bullets(s, Inches(0.95), Inches(2.7), Inches(11.4), Inches(3.5),
        ["End-to-end system: data pipeline → two models → web app",
         "Custom CNN and ResNet50 trained, evaluated & auto-compared",
         "Best model selected via a weighted composite score",
         "Detects 7 lung conditions from chest X-rays",
         "Designed as a decision-support aid for clinicians"],
        size=18, color=RGBColor(0xDD,0xE6,0xF2), gap=14)
text(s, Inches(0.95), Inches(6.5), Inches(11.5), Inches(0.6),
     [[("Thank you  —  Questions?", 22, ACCENT, True)]])

# ── Save ─────────────────────────────────────────────────────────────────────
out = Path(__file__).parent / "LungAI_Presentation.pptx"
prs.save(str(out))
print(f"Saved: {out}  ({len(prs.slides._sldIdLst)} slides)")
